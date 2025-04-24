# ChatBot_BE/document_qa/doc_parser.py

import os
import docx
import pptx
from pdfminer.high_level import extract_text

def extract_text_from_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(path):
    prs = pptx.Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(path):
    return extract_text(path)

def parse_all_documents(folder_path="documents"):
    parsed_docs = []

    for fname in os.listdir(folder_path):
        fpath = os.path.join(folder_path, fname)
        if not os.path.isfile(fpath):
            continue
        try:
            if fname.endswith(".pdf"):
                text = extract_text_from_pdf(fpath)
            elif fname.endswith(".docx"):
                text = extract_text_from_docx(fpath)
            elif fname.endswith(".pptx"):
                text = extract_text_from_pptx(fpath)
            else:
                continue  # skip unsupported
            parsed_docs.append({
                "filename": fname,
                "text": text.strip()
            })
        except Exception as e:
            print(f"[Error] {fname}: {str(e)}")
    
    return parsed_docs
